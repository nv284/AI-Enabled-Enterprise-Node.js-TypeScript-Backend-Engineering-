from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Helper function to add a slide with bullet points
    def add_content_slide(title_text, content_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = content_list[0]
        for item in content_list[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Advanced API Security & Auth Architectures"
    subtitle.text = "JWT, RBAC, OAuth2/OIDC, Secure API Design & AI Audits\nTarget Audience: Expert TypeScript Engineers"

    # Slide 2: Agenda
    add_content_slide("Agenda", [
        "1. JWT: Advanced Implementation & Pitfalls",
        "2. RBAC: Architecture & TS Decorators",
        "3. OAuth2 & OIDC: PKCE & Token Management",
        "4. Secure API Design: OWASP Top 10",
        "5. API Security Implementation in TypeScript",
        "6. AI-Assisted Security Audits & Vulnerability Detection",
        "7. AI Tools in the TS Ecosystem",
        "8. Summary & Best Practices"
    ])

    # Slide 3: JWT
    add_content_slide("1. JWT: Advanced Implementation & Pitfalls", [
        "• Library Choice: Use 'jose' instead of 'jsonwebtoken' (WebCrypto, non-blocking, Edge-ready).",
        "• Algorithm Confusion: Strictly enforce 'alg' (e.g., RS256). Never trust client-provided 'alg'.",
        "• Stateless Revocation: Use short-lived Access Tokens (15m) + secure HttpOnly Refresh Tokens.",
        "• TS Implementation: import * as jwt from 'jose'; const { payload } = await jwt.jwtVerify(token, publicKey);"
    ])

    # Slide 4: RBAC
    add_content_slide("2. RBAC: Architecture & TS Decorators", [
        "• Beyond Simple Arrays: Move to hierarchical roles or ABAC for complex domains.",
        "• TS Decorators: Use NestJS-style decorators for clean, declarative authorization.",
        "• Policy Enforcement: Centralize logic using engines like Casbin.",
        "• TS Example: @Roles('ADMIN', 'MANAGER') @UseGuards(RolesGuard) async deleteUser(...) { ... }"
    ])

    # Slide 5: OAuth2/OIDC
    add_content_slide("3. OAuth2 & OIDC: PKCE & Token Management", [
        "• SPA / Mobile Clients: MUST use Authorization Code Flow with PKCE (RFC 7636).",
        "• OIDC Scopes: Use 'openid profile email' carefully. Request least privilege.",
        "• Token Validation: Never decode JWTs manually. Always verify signature against JWKS.",
        "• Back-channel logout: Implement OIDC Back-Channel Logout for immediate session termination."
    ])

    # Slide 6: Secure API Design
    add_content_slide("4. Secure API Design: OWASP API Top 10", [
        "• BOLA / IDOR (API1): Always verify resource ownership in the DB query, not just in memory.",
        "• Mass Assignment (API8): Never bind raw request bodies to DB models. Use strict DTOs.",
        "• Rate Limiting: Implement per-user, per-IP, and per-endpoint rate limiting.",
        "• Pagination & Filtering: Enforce strict limits to prevent NoSQL/SQL injection and resource exhaustion."
    ])

    # Slide 7: API Security in TS
    add_content_slide("5. API Security Implementation in TypeScript", [
        "• Strict Validation: Use 'zod' or 'class-validator' to enforce DTO schemas at the edge.",
        "• Security Headers: Use 'helmet' to set CSP, HSTS, X-Frame-Options.",
        "• Error Handling: Never leak stack traces or internal DB schemas in 500 errors.",
        "• TS Example (Zod): const UserSchema = z.object({ id: z.string().uuid() }); const safeData = UserSchema.parse(req.body);"
    ])

    # Slide 8: AI-Assisted Security Audits
    add_content_slide("6. AI-Assisted Security Audits & Vulnerability Detection", [
        "• LLM-Powered SAST: AI models analyze AST to find context-aware vulnerabilities that regex misses.",
        "• Automated Threat Modeling: AI generates Data Flow Diagrams (DFD) and suggests mitigations.",
        "• Smart Dependency Analysis: AI evaluates transitive dependencies for reachability.",
        "• Remediation: AI suggests exact code patches (e.g., GitHub Copilot Workspace, Snyk DeepCode)."
    ])

    # Slide 9: AI Tools in TS
    add_content_slide("7. AI Tools in the TypeScript Ecosystem", [
        "• Snyk & Socket: AI-driven analysis of npm packages for malware, typosquatting, and vulnerable code paths.",
        "• GitHub Advanced Security: CodeQL combined with AI for semantic code analysis in TS/JS.",
        "• ESLint Security Plugins: 'eslint-plugin-security' enhanced with AI rulesets for custom patterns.",
        "• Continuous Monitoring: AI agents monitor production logs for anomalous API usage patterns."
    ])

    # Slide 10: Summary
    add_content_slide("8. Summary & Best Practices", [
        "• Zero Trust Architecture: Never trust the client. Validate every JWT, request body, and DB query.",
        "• TypeScript as a Security Tool: Leverage strict typing to eliminate Mass Assignment and Injection.",
        "• Shift Left with AI: Integrate AI security audits directly into the CI/CD pipeline.",
        "• Defense in Depth: Combine RBAC, Secure API Design, and continuous AI monitoring."
    ])

    # Save the presentation
    prs.save('Advanced_API_Security_TS.pptx')
    print("Presentation saved successfully as 'Advanced_API_Security_TS.pptx'")

if __name__ == "__main__":
    create_presentation()