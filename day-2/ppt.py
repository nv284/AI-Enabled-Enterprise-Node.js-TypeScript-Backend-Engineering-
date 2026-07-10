from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat
import requests
from io import BytesIO

# Create presentation
prs = Presentation()

# Helper function to add image to slide
def add_image_to_slide(slide, image_url, left, top, width=None, height=None):
    try:
        response = requests.get(image_url)
        img_stream = BytesIO(response.content)
        slide.shapes.add_picture(img_stream, left, top, width, height)
    except:
        # Add placeholder text if image fails
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "[Image: " + image_url.split('/')[-1] + "]"
        tf.paragraphs[0].font.size = Pt(10)

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Modern API Development:\nContract-First Approach"
subtitle.text = "A Guide for Senior Developers\nUnderstanding Modern API Design Tools & Practices"

# Slide 2: Evolution of API Development
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Evolution of API Development"
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "Traditional Development (2000-2010)"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "• Monolithic applications"
p.level = 1
p = text_frame.add_paragraph()
p.text = "• No formal API contracts needed"
p.level = 1

p = text_frame.add_paragraph()
p.text = "Service-Oriented Architecture (SOA)"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "• SOAP/XML-based services"
p.level = 1
p = text_frame.add_paragraph()
p.text = "• WSDL as contract (complex & cumbersome)"
p.level = 1

p = text_frame.add_paragraph()
p.text = "REST APIs (2010-2018)"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "• Code-First approach dominant"
p.level = 1
p = text_frame.add_paragraph()
p.text = "• Documentation often an afterthought"
p.level = 1

p = text_frame.add_paragraph()
p.text = "Modern Era (2018-Present)"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "• Contract-First development"
p.level = 1
p = text_frame.add_paragraph()
p.text = "• OpenAPI specification as source of truth"
p.level = 1

# Slide 3: Problems with Code-First Development
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Problems with Code-First Development"
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "❌ Inconsistent API Design"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "Different developers create different endpoint styles"
p.level = 1

p = text_frame.add_paragraph()
p.text = "❌ Documentation Becomes Outdated"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "Code changes but documentation doesn't"
p.level = 1

p = text_frame.add_paragraph()
p.text = "❌ Teams Block Each Other"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "Frontend, QA, and Mobile teams wait for backend completion"
p.level = 1

p = text_frame.add_paragraph()
p.text = "❌ Different Response Formats"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "{books: []}, {data: []}, {result: []} - lack of consistency"
p.level = 1

p = text_frame.add_paragraph()
p.text = "❌ Late Discovery of Design Issues"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "Expensive refactoring required"
p.level = 1

# Slide 4: What is Contract-First Development?
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "What is Contract-First Development?"
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "📋 Definition"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "Design and agree on the API specification BEFORE writing implementation"
p.level = 0

p = text_frame.add_paragraph()
p.text = " Key Principles"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "OpenAPI document becomes the contract between producers and consumers"
p.level = 1
p = text_frame.add_paragraph()
p.text = "Everyone agrees on endpoints, request/response formats first"
p.level = 1
p = text_frame.add_paragraph()
p.text = "Only then does coding begin"
p.level = 1

p = text_frame.add_paragraph()
p.text = "📝 Contract Covers"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "What endpoints exist?"
p.level = 1
p = text_frame.add_paragraph()
p.text = "What request body is expected?"
p.level = 1
p = text_frame.add_paragraph()
p.text = "What does the response look like?"
p.level = 1
p = text_frame.add_paragraph()
p.text = "Error formats & validation rules"
p.level = 1

# Slide 5: Traditional vs Contract-First Workflow
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Traditional vs Contract-First Workflow"
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "🔄 Traditional Workflow (Sequential)"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "1. Requirement → 2. Developer writes code → 3. Backend completed"
p.level = 1
p = text_frame.add_paragraph()
p.text = "4. Frontend starts → 5. QA starts → 6. Documentation written"
p.level = 1
p = text_frame.add_paragraph()
p.text = "Result: Teams wait, delays compound"
p.level = 1

p = text_frame.add_paragraph()
p.text = "⚡ Contract-First Workflow (Parallel)"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "1. Requirement → 2. Design OpenAPI Contract → 3. Review & Approve"
p.level = 1
p = text_frame.add_paragraph()
p.text = "4. Frontend Starts + Backend Starts + QA Starts (SIMULTANEOUSLY)"
p.level = 1
p = text_frame.add_paragraph()
p.text = "5. Documentation Ready → 6. Deployment"
p.level = 1
p = text_frame.add_paragraph()
p.text = "Result: Parallel development, faster delivery"
p.level = 1

# Slide 6: Key Benefits
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Key Benefits of Contract-First Development"
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

benefits = [
    ("Single Source of Truth", "Everyone refers to the same OpenAPI document"),
    ("Parallel Development", "Frontend, Backend, QA work simultaneously"),
    ("Automatic Documentation", "Generated from contract using Swagger UI"),
    ("Strong Typing", "Generate TypeScript, Java, C#, Python models automatically"),
    ("Early Validation", "Review design before coding - cheaper fixes"),
    ("Better Collaboration", "All teams work from same specification")
]

for benefit, description in benefits:
    p = text_frame.add_paragraph()
    p.text = f"✅ {benefit}"
    p.font.bold = True
    p = text_frame.add_paragraph()
    p.text = description
    p.level = 1
    p.font.size = Pt(14)

# Slide 7: Real-World Example
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Real-World Example: E-Commerce Platform"
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "🛒 Scenario"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "Multiple clients depend on the same API:"
p.level = 0
clients = ["Web Application", "Mobile App", "Warehouse System", "Payment Gateway", "Analytics Platform", "Customer Support Portal"]
for client in clients:
    p = text_frame.add_paragraph()
    p.text = f"• {client}"
    p.level = 1

p = text_frame.add_paragraph()
p.text = "⚠️ Problem Without Contract"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "Backend changes 'price' field to 'productPrice' → All clients break!"
p.level = 0

p = text_frame.add_paragraph()
p.text = "✅ Solution With Contract-First"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "1. Change proposed in OpenAPI spec → 2. Teams review & approve"
p.level = 1
p = text_frame.add_paragraph()
p.text = "3. Updated types/SDKs generated → 4. Consumers adapt before deployment"
p.level = 1

# Slide 8: Why Necessary for Modern Applications
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Why Contract-First for Modern Applications?"
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "🌐 Modern Systems Are:"
p.font.bold = True
modern_traits = [
    "Distributed & Microservice-based",
    "Cloud-native",
    "API-driven",
    "Consumed by many different clients",
    "Integrated with partners & AI systems"
]
for trait in modern_traits:
    p = text_frame.add_paragraph()
    p.text = f"• {trait}"
    p.level = 1

p = text_frame.add_paragraph()
p.text = "\n💡 Key Insight"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "APIs are PRODUCTS, not just implementation details"
p.level = 0

p = text_frame.add_paragraph()
p.text = "\n🎯 Result"
p.font.bold = True
p = text_frame.add_paragraph()
p.text = "Clear, versioned contract provides stability and enables independent teams to work efficiently"
p.level = 0

# Slide 9: Comparison Table
slide_layout = prs.slide_layouts[5]  # Title Only
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Traditional vs Contract-First: Side-by-Side"

# Create table
rows = 7
cols = 3
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9.0)
height = Inches(0.8)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(3.0)
table.columns[2].width = Inches(3.0)

# Headers
table.cell(0, 0).text = "Aspect"
table.cell(0, 1).text = "Traditional (Code-First)"
table.cell(0, 2).text = "Contract-First"

# Data
data = [
    ("Source of Truth", "Code", "OpenAPI Contract"),
    ("Documentation", "Often lags behind", "Generated from contract"),
    ("Development", "Frontend waits for backend", "Parallel development"),
    ("Model Creation", "Manual", "Auto-generated"),
    ("Design Issues", "Late discovery", "Reviewed before coding"),
    ("Integration Risk", "Higher", "Consistent across teams")
]

for i, (aspect, traditional, contract) in enumerate(data, 1):
    table.cell(i, 0).text = aspect
    table.cell(i, 1).text = traditional
    table.cell(i, 2).text = contract

# Slide 10: Closing & Key Takeaway
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Key Takeaway"
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "🏗️ The Blueprint Analogy"
p.font.bold = True
p.font.size = Pt(18)

p = text_frame.add_paragraph()
p.text = "\n\"Think of the OpenAPI specification as a blueprint for building a house. You don't start laying bricks and then decide where the doors and windows go—you first agree on the architectural plan.\""
p.font.size = Pt(14)
p.font.italic = True

p = text_frame.add_paragraph()
p.text = "\n📋 Contract-First Benefits:"
p.font.bold = True
p.font.size = Pt(16)

benefits_list = [
    "All teams aligned from the start",
    "Reduced integration issues",
    "Parallel development enabled",
    "Automatic documentation & type generation",
    "Early validation of design decisions"
]

for benefit in benefits_list:
    p = text_frame.add_paragraph()
    p.text = f"✓ {benefit}"
    p.level = 1
    p.font.size = Pt(14)

p = text_frame.add_paragraph()
p.text = "\n🎯 Contract-First is now the standard for modern API-driven applications!"
p.font.bold = True
p.font.size = Pt(16)

# Save presentation
prs.save('API_Design_Tools_Presentation.pptx')
print("Presentation created successfully: API_Design_Tools_Presentation.pptx")